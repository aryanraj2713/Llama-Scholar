from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
import uvicorn
import os
import PyPDF2
from pptx import Presentation
import io
import boto3
from botocore.exceptions import ClientError
from PIL import Image
import tempfile
import base64
from groq import Groq
import json
import requests
from pydantic import BaseModel
import httpx
import uuid
import datetime
import pytz
# from typing import List, Dict, Any
# from qdrant_client import QdrantClient
# from qdrant_client.http import models
# import httpx
# from sentence_transformers import SentenceTransformer
# import nltk
# from nltk.tokenize import sent_tokenize

app = FastAPI()

class Task(BaseModel):
    epoch_time: int
    name: str
    email: str
    task_name: str

# AWS Bedrock configuration
aws_access_key_id = "AKIAQ3EGUNCQ5QJVNYMJ"
aws_secret_access_key = "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"
aws_region = "us-west-2"

TUNE_AI_API_URL = "https://proxy.tune.app/chat/completions"
TUNE_AI_API_KEY = "sk-tune-IeArLOH0xZIE5cTHuNzrVGQlATFTMNDZqzh"  # Replace with your actual API key

bedrock_client = boto3.client(
    service_name="bedrock-runtime",
    region_name=aws_region,
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key
)

# Groq configuration
groq_client = Groq(
    api_key="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",
)
# Download NLTK data
# nltk.download('punkt')



# # Initialize Qdrant client with cloud configuration
# QDRANT_URL = "https://13f743e4-bf99-4752-a59a-fb717c1b8d52.us-east4-0.gcp.cloud.qdrant.io:6333"
# QDRANT_API_KEY = os.environ.get("QDRANT_API_KEY", "<your-token>")  # Replace with your actual Qdrant API key

# qdrant_client = QdrantClient(
#     url=QDRANT_URL,
#     api_key=QDRANT_API_KEY,
# )

# # Initialize sentence transformer model
# model = SentenceTransformer('all-MiniLM-L6-v2')

@app.get("/")
async def root():
    return {"message": "Welcome"}

@app.get("/greeting")
async def greeting():
    return {"message": "Hello", "details": "You are using LAamAScholar Bot"}

@app.post("/summarize")
async def summarize_file(file: UploadFile = File(...)):
    content = await file.read()
    
    if file.filename.endswith('.pdf'):
        text = extract_text_from_pdf(content)
    elif file.filename.endswith('.pptx'):
        text = extract_text_from_pptx(content)
    else:
        return JSONResponse(status_code=400, content={"error": "Unsupported file format. Please upload a PDF or PPTX file."})
    
    summary = summarize_text(text)
    return {"summary": summary}

@app.post("/generate_questions")
async def generate_questions(file: UploadFile = File(...)):
    content = await file.read()
    
    if file.filename.endswith('.pdf'):
        text = extract_text_from_pdf(content)
    elif file.filename.endswith('.pptx'):
        text = extract_text_from_pptx(content)
    else:
        return JSONResponse(status_code=400, content={"error": "Unsupported file format. Please upload a PDF or PPTX file."})
    
    questions = generate_important_questions(text)
    return {"questions": questions}

@app.post("/voice-notes")
async def transcribe_voice_notes(file: UploadFile = File(...)):
    if not file.filename.endswith('.mp3'):
        return JSONResponse(status_code=400, content={"error": "Unsupported file format. Please upload an MP3 file."})
    
    # Create a temporary file to store the uploaded content
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_file:
        content = await file.read()
        temp_file.write(content)
        temp_file_path = temp_file.name

    try:
        with open(temp_file_path, "rb") as audio_file:
            transcription = groq_client.audio.transcriptions.create(
                file=(file.filename, audio_file),
                model="whisper-large-v3",
                response_format="verbose_json",
            )
        
        return {"transcription": transcription.text}
    finally:
        # Clean up the temporary file
        os.unlink(temp_file_path)

@app.post("/ocr/")
async def perform_ocr(file: UploadFile = File(...)):
    # Read and encode the image
    contents = await file.read()
    base64_image = encode_image(io.BytesIO(contents))

    # Prepare the message for Tune AI
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{base64_image}"
                    }
                },
                {
                    "type": "text",
                    "text": "Explain what is written in the image from an educational Point of View."
                }
            ]
        }
    ]

    # Prepare the request payload
    payload = {
        "temperature": 0.8,
        "messages": messages,
        "model": "meta/llama-3.2-90b-vision",
        "stream": False,
        "frequency_penalty": 0,
        "max_tokens": 900
    }

    # Make the API call to Tune AI
    headers = {
        "Authorization": f"Bearer {TUNE_AI_API_KEY}",
        "Content-Type": "application/json"
    }

    response = requests.post(TUNE_AI_API_URL, headers=headers, json=payload)
    response.raise_for_status()
    result = response.json()

    # Extract the OCR result from the response
    ocr_result = result['choices'][0]['message']['content']

    return {"ocr_result": ocr_result}

# Initialize sentence transformer model
# model = SentenceTransformer('all-MiniLM-L6-v2')

# # Define the collection name
# COLLECTION_NAME = "document_collection"

# # Define request models
# class SearchQuery(BaseModel):
#     query: str

# # Function to create the collection
# def create_collection():
#     qdrant_client.recreate_collection(
#         collection_name=COLLECTION_NAME,
#         vectors_config=models.VectorParams(
#             size=model.get_sentence_embedding_dimension(),
#             distance=models.Distance.COSINE,
#         ),
#     )

# # Function to chunk text
# def chunk_text(text: str, chunk_size: int = 5):
#     sentences = sent_tokenize(text)
#     chunks = []
#     current_chunk = []
#     current_length = 0
    
#     for sentence in sentences:
#         current_chunk.append(sentence)
#         current_length += 1
#         if current_length >= chunk_size:
#             chunks.append(" ".join(current_chunk))
#             current_chunk = []
#             current_length = 0
    
#     if current_chunk:
#         chunks.append(" ".join(current_chunk))
    
#     return chunks

# # Function to process and insert document
# def process_and_insert_document(file_content: str, file_name: str):
#     chunks = chunk_text(file_content)
#     embeddings = model.encode(chunks)
    
#     points = []
#     for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
#         points.append(
#             models.PointStruct(
#                 id=f"{file_name}_{i}",
#                 vector=embedding.tolist(),
#                 payload={
#                     "text": chunk,
#                     "file_name": file_name,
#                     "chunk_index": i
#                 }
#             )
#         )
    
#     qdrant_client.upsert(
#         collection_name=COLLECTION_NAME,
#         points=points
#     )

# # Initialize the collection on startup
# @app.on_event("startup")
# async def startup_event():
#     # Check if the collection exists, create it if it doesn't
#     collections = qdrant_client.get_collections()
#     if COLLECTION_NAME not in [collection.name for collection in collections.collections]:
#         create_collection()
#     print(f"Connected to Qdrant. Available collections: {qdrant_client.get_collections().collections}")

# @app.post("/upload")
# async def upload_file(file: UploadFile = File(...), api_key: APIKey = Depends(get_api_key)):
#     content = await file.read()
#     file_content = content.decode("utf-8")
#     process_and_insert_document(file_content, file.filename)
#     return {"message": f"File {file.filename} uploaded and processed successfully"}

# @app.post("/docsearch")
# async def docsearch(query: SearchQuery, api_key: APIKey = Depends(get_api_key)):
#     try:
#         query_vector = model.encode([query.query])[0]
#         search_result = qdrant_client.search(
#             collection_name=COLLECTION_NAME,
#             query_vector=query_vector.tolist(),
#             limit=5
#         )
        
#         return [
#             {
#                 "id": result.id,
#                 "score": result.score,
#                 "text": result.payload["text"],
#                 "file_name": result.payload["file_name"],
#                 "chunk_index": result.payload["chunk_index"]
#             } for result in search_result
#         ]
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=str(e))

def extract_text_from_pdf(content):
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def summarize_text(text):
    prompt = f"Summarize the following text:\n\n{text[:4000]}..."  # Truncate to 4000 characters to fit within token limit
    
    response = bedrock_client.invoke_model(
        modelId="meta.llama3-1-70b-instruct-v1:0",
        body=json.dumps({
            "prompt": prompt
        })
    )
    
    response_body = json.loads(response['body'].read())
    return response_body['generation']

def generate_important_questions(text):
    prompt = f"Generate 5 important questions based on the following text:\n\n{text[:4000]}..."  # Truncate to 4000 characters to fit within token limit
    
    response = bedrock_client.invoke_model(
        modelId="meta.llama3-1-70b-instruct-v1:0",
        body=json.dumps({
            "prompt": prompt
        })
    )
    
    response_body = json.loads(response['body'].read())
    return response_body['generation']


def encode_image(image_file):
    return base64.b64encode(image_file.read()).decode('utf-8')

'''
    task service for storing task and sending out as reminders
'''
SUPABASE_URL = "https://jpoqergdaandvpxyirtq.supabase.co"
SUPABASE_API_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Impwb3FlcmdkYWFuZHZweHlpcnRxIiwicm9sZSI6ImFub24iLCJpYXQiOjE3MjkzMTc0NjIsImV4cCI6MjA0NDg5MzQ2Mn0.OZTinuAeJObdhIv9AfzMXWZcs0aofaUAn611HpEKWfs"

@app.post("/tasks")
async def create_task(task: Task):
    # Prepare the data to send to Supabase
    data = {
        "epoch_time": task.epoch_time,
        "name": task.name,
        "email": task.email,
        "task_name": task.task_name,
    }

    # Send a request to Supabase
    async with httpx.AsyncClient() as client:
        response = await client.post(
            f"{SUPABASE_URL}/rest/v1/tasks",
            json=data,
            headers={
                "Authorization": f"Bearer {SUPABASE_API_KEY}",
                "Content-Type": "application/json",
                "apikey": SUPABASE_API_KEY,
            },
        )
    create_eventbridge_schedule(data)
    # Check if the request was successful
    if response.status_code == 201:
        return {"message": "Task created successfully", "data": data}
    else:
        raise HTTPException(status_code=response.status_code, detail=response.json())

@app.get("/tasks/{email}")
async def get_tasks(email: str):
    # Send a request to Supabase
    async with httpx.AsyncClient() as client:
        response = await client.get(
            f"{SUPABASE_URL}/rest/v1/tasks?select=*&email=eq.{email}",
            headers={
                "Authorization": f"Bearer {SUPABASE_API_KEY}",
                "Content-Type": "application/json",
                "apikey": SUPABASE_API_KEY,
            },
        )

    # Check if the request was successful
    if response.status_code == 200:
        return response.json()
    else:
        raise HTTPException(status_code=response.status_code, detail=response.json())

def create_eventbridge_schedule(task):
    # Initialize the Boto3 clients for EventBridge and Lambda
    lambda_client = boto3.client(
        'lambda',
        aws_access_key_id='AKIAR76O52QWHA2RBJ7U',
        aws_secret_access_key='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',
        region_name='us-west-2')
    eventbridge_client = boto3.client(
        'scheduler',
        aws_access_key_id='AKIAR76O52QWHA2RBJ7U',
        aws_secret_access_key='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',
        region_name='us-west-2')

    # Generate a unique ID for the rule
    schedule_id = str(uuid.uuid4())
    schedule_name = f"schedule-{schedule_id}"
    lambda_function_arn = "arn:aws:lambda:us-west-2:137336050732:function:meta-scheduler-lambda"

    eventbridge_client.create_schedule(
        Name=schedule_name,
        ScheduleExpression=f"cron({convert_epoch_to_cron(task['epoch_time'])})",
        Target={
            'Arn': lambda_function_arn,
            'Input': json.dumps(task),
            'RoleArn': 'arn:aws:iam::137336050732:role/service-role/meta-scheduler-lambda-role-omq2exfb'
        },
        State='ENABLED',
        Description='My scheduled event',
        FlexibleTimeWindow={
        'Mode': 'OFF',  # Set Mode to 'OFF' if you don't want flexible time windows
        # 'MaximumWindowInMinutes': 300  # Use this if flexible time window is required
        },
        StartDate=datetime.datetime.now(
            pytz.timezone("Asia/Kolkata")),  # Start date of the schedule (optional, default is the current date and time)
        EndDate=datetime.datetime.now(
            pytz.timezone("Asia/Kolkata")) + datetime.timedelta(days=1)  # End date of the schedule (optional, default is 1 day from the current date and time)
            
    )

    print(f"Created schedule with ID: {schedule_id}")

def convert_epoch_to_cron(epoch_time):
    # Convert epoch time to a datetime object
    dt = datetime.datetime.fromtimestamp(epoch_time, tz=pytz.timezone("Asia/Kolkata"))
    
    # Ensure the cron expression is for a future time
    current_time = datetime.datetime.now(pytz.timezone("Asia/Kolkata"))
    
    # If the task time is in the past, adjust it to a future time (e.g., 5 minutes from now)
    if dt <= current_time:
        dt = current_time + datetime.timedelta(minutes=5)
    
    # Create the cron expression (minute hour day-of-month month day-of-week year)
    cron_expression = f"{dt.minute} {dt.hour} {dt.day} {dt.month} ? {dt.year}"
    
    return cron_expression

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)