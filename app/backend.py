from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from typing import Optional
from fastapi.responses import JSONResponse, PlainTextResponse
import uvicorn

import PyPDF2
from pptx import Presentation
import io
import boto3
from botocore.exceptions import ClientError

import json

from pydantic import BaseModel
import httpx
from qdrant_client import QdrantClient
from qdrant_client.http import models
import uuid
import logging


logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                    handlers=[logging.StreamHandler()])

logger = logging.getLogger("app")
app = FastAPI()



# Clients
class MemoryItem(BaseModel):
    content: str
    metadata: dict = {}

class Question(BaseModel):
    text: str
# AWS Bedrock configuration
aws_access_key_id = "AKIAQ3EGUNCQ5QJVNYMJ"
aws_secret_access_key = "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"
aws_region = "us-west-2"

TUNE_AI_API_URL = "https://proxy.tune.app/chat/completions"
TUNE_AI_API_KEY = "sk-tune-IeArLOH0xZIE5cTHuNzrVGQlATFTMNDZqzh" 

qdrant_client = QdrantClient(
    url="https://13f743e4-bf99-4752-a59a-fb717c1b8d52.us-east4-0.gcp.cloud.qdrant.io:6333", 
    api_key="16WN21MOBd-WI1JNKD7KTQ8Pn378xUv9PY4MzQQiwGNPf3ZUf3gZhg",
)


bedrock_client = boto3.client(
    service_name="bedrock-runtime",
    region_name=aws_region,
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key
)








#routes
@app.get("/")
async def root():
    return {"message": "Welcome"}

@app.get("/greeting")
async def greeting():
    return {"message": "Hello", "details": "You are using LAamAScholar Bot"}

@app.post("/summarize")
async def summarize_file(file: UploadFile = File(...)):
    try:
        content = await file.read()
        if file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(content)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(content)
        else:
            return JSONResponse(status_code=400, content={"error": "Unsupported file format. Please upload a PDF or PPTX file."})
        
        summary = summarize_text(text)
        
        if summary.startswith("Error:"):
            return JSONResponse(status_code=500, content={"error": summary})
        
        return {"summary": summary}
    except Exception as e:
        logger.error(f"Unexpected error in summarize_file: {str(e)}")
        return JSONResponse(status_code=500, content={"error": f"An unexpected error occurred: {str(e)}"})

@app.post("/remember")
async def remember(
    file: UploadFile = File(...),
    metadata: Optional[str] = Form(None)
):
    try:
        logger.info("Processing remember request")
        
        # Generate a unique ID for this memory
        memory_id = str(uuid.uuid4())
        logger.info(f"Generated memory ID: {memory_id}")
        
        # Parse metadata
        metadata_dict = {}
        if metadata:
            if metadata.lower() == 'string' or metadata.strip() == '':
                logger.info("Metadata is empty or 'string', using empty dict")
            else:
                try:
                    metadata_dict = json.loads(metadata)
                    logger.info(f"Parsed metadata: {metadata_dict}")
                except json.JSONDecodeError as e:
                    logger.warning(f"Could not parse metadata as JSON: {str(e)}. Using as string value.")
                    metadata_dict = {"value": metadata}
        
        # Read and extract content from the file
        content = await file.read()
        if file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(content)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(content)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format. Please upload a PDF or PPTX file.")
        
        logger.info(f"Extracted text (first 100 chars): {text[:100]}...")
        
        # Generate embedding using AWS Bedrock Titan model
        logger.info("Generating embedding")
        try:
            embedding_response = bedrock_client.invoke_model(
                modelId="amazon.titan-embed-text-v1",
                body=json.dumps({
                    "inputText": text
                })
            )
            embedding_body = json.loads(embedding_response['body'].read())
            if 'embedding' not in embedding_body:
                logger.error(f"Unexpected response from Bedrock: {embedding_body}")
                raise HTTPException(status_code=500, detail="Unexpected response from Bedrock API")
            embedding = embedding_body['embedding']
            logger.info(f"Embedding generated successfully. Length: {len(embedding)}")
        except Exception as e:
            logger.error(f"Error generating embedding: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error generating embedding: {str(e)}")
        
        # Save to Qdrant
        logger.info("Saving to Qdrant")
        try:
            qdrant_client.upsert(
                collection_name="memories",
                points=[
                    models.PointStruct(
                        id=memory_id,
                        vector=embedding,
                        payload={
                            "content": text,
                            "filename": file.filename,
                            **metadata_dict
                        }
                    )
                ]
            )
            logger.info("Successfully saved to Qdrant")
        except Exception as e:
            logger.error(f"Error saving to Qdrant: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error saving to Qdrant: {str(e)}")
        
        return {"message": "Memory saved successfully", "id": memory_id}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in remember endpoint: {str(e)}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")

@app.post("/response-vector")
async def response_vector(question: Question):
    try:
        # Generate embedding for the question
        try:
            question_embedding = generate_embedding(question.text)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating question embedding: {str(e)}")

        # Search for similar vectors in Qdrant
        try:
            search_results = qdrant_client.search(
                collection_name="memories",
                query_vector=question_embedding,
                limit=3  # Adjust as needed
            )
        except Exception as e:
            logger.exception(f"Error searching Qdrant: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error searching vector database: {str(e)}")

        # Retrieve and compile related content
        context = ""
        for result in search_results:
            context += result.payload.get('content', '') + "\n\n"

        # Generate response using Bedrock
        prompt = f"""Context: {context}

        Question: {question.text}

        Please provide a concise and accurate answer to the question based on the given context. If the context doesn't contain relevant information to answer the question, please state that you don't have enough information to provide an accurate answer.

        Answer:"""

        try:
            response = bedrock_client.invoke_model(
                modelId="meta.llama3-1-70b-instruct-v1:0",
                body=json.dumps({
                    "prompt": prompt,
                    "temperature": 0.7,
                    "top_p": 0.95,
                })
            )
            response_body = json.loads(response['body'].read())
            answer = response_body.get('generation', '').strip()
            if not answer:
                logger.error(f"Unexpected response from Bedrock: {response_body}")
                raise ValueError("Unexpected response format from Bedrock API")
        except Exception as e:
            logger.exception(f"Error generating answer: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error generating answer: {str(e)}")

        return {"answer": answer}

    except HTTPException:
        raise
    except Exception as e:
        logger.exception(f"Unexpected error in response_vector endpoint: {str(e)}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")


@app.post("/generate_questions")
async def generate_questions(file: UploadFile = File(...)):
    content = await file.read()
    
    if file.filename.endswith('.pdf'):
        text = extract_text_from_pdf(content)
    elif file.filename.endswith('.pptx'):
        text = extract_text_from_pptx(content)
    else:
        return JSONResponse(status_code=400, content={"error": "Unsupported file format. Please upload a PDF or PPTX file."})
    
    questions = generate_important_questions(text)
    return {"questions": questions}
@app.get("/webhook")
async def webhook(request: Request):
    # Handle the webhook verification
    mode = request.query_params.get('hub.mode')
    token = request.query_params.get('hub.verify_token')
    challenge = request.query_params.get('hub.challenge')

    if mode and token:
        if mode == 'subscribe' and token == 'meatyhamhock':
            print('WEBHOOK_VERIFIED')
            return PlainTextResponse(content=challenge)
        else:
            raise HTTPException(status_code=403, detail="Verification failed")
    
    return PlainTextResponse(content="OK")


@app.post("/webhook")
async def postWebhook(request: Request):
    
    data = await request.json() 
    logger.info(data)

    if(data['entry'][0]['changes'][0]['value']['statuses'] is not None):
        return "OK"


    url = "https://graph.facebook.com/v20.0/436207996245933/messages"
    headers = {
        "Authorization": "Bearer EAAPZCdgo11ucBOxmxq3h4XlhzL7IituEG4L5yXnQbMiD0LtyakeLD0bZB0ZBdczo18hIJhW6ZAQ8aPRkw127wlXnD7Cks7IWO3KELZBHfa1jIwGZAiX9dL6ftEk4Oxng3gVdZCRgG8QlGtsEH0Hcf4j62ZAKOLZALdcuZCZC4HIp1ZBhFi65hzr2ACuGnAqCGZAS0abrfQP7KHydrvKMSwAPHfZBqnXUtZBZAHcZD",
        "Content-Type": "application/json"
    }
    payload = {
        "messaging_product": "whatsapp",
        "to": "919554469077",
        "type": "template",
        "template": {
            "name": "hello_world",
            "language": {
                "code": "en_US"
            }
        }
    }

    async with httpx.AsyncClient() as client:
        response = await client.post(url, headers=headers, json=payload)


    if response.status_code == 200:
        return {"message": "WhatsApp message sent successfully"}
    else:
        raise HTTPException(status_code=response.status_code, detail="Failed to send WhatsApp message")



















# Functions for extracting text from PDF and PPTX files
def extract_text_from_pdf(content):
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def summarize_text(text):
    prompt = f"Summarize the following text:\n\n{text[:4000]}..."  
    try:
        response = bedrock_client.invoke_model(
            modelId="meta.llama3-1-70b-instruct-v1:0",
            body=json.dumps({
                "prompt": prompt
            })
        )
        
        response_body = json.loads(response['body'].read())
        logger.info(f"Bedrock API response: {response_body}")
        
        if 'generation' in response_body:
            return response_body['generation'].strip()
        else:
            logger.error(f"Unexpected response format: {response_body}")
            return "Error: Unexpected response format from Bedrock API"
    
    except ClientError as e:
        logger.error(f"AWS Bedrock ClientError: {str(e)}")
        return f"Error: AWS Bedrock API error - {str(e)}"
    except json.JSONDecodeError as e:
        logger.error(f"JSON Decode Error: {str(e)}")
        return "Error: Failed to parse Bedrock API response"
    except Exception as e:
        logger.error(f"Unexpected error in summarize_text: {str(e)}")
        return f"Error: Unexpected error occurred - {str(e)}"

def generate_embedding(text):
    try:
        logger.info("Generating embedding")
        embedding_response = bedrock_client.invoke_model(
            modelId="amazon.titan-embed-text-v1",
            body=json.dumps({
                "inputText": text
            })
        )
        
        logger.debug(f"Raw Bedrock response: {embedding_response}")
        
        embedding_body = json.loads(embedding_response['body'].read())
        logger.debug(f"Parsed embedding body: {embedding_body}")
        
        if 'embedding' not in embedding_body:
            logger.error(f"'embedding' key not found in response. Full response: {embedding_body}")
            raise ValueError("Unexpected response format from Bedrock API")
        
        embedding = embedding_body['embedding']
        logger.info(f"Embedding generated successfully. Length: {len(embedding)}")
        return embedding
    except Exception as e:
        logger.exception(f"Error generating embedding: {str(e)}")
        raise

def generate_important_questions(text):
    prompt = f"Generate 5 important questions based on the following text:\n\n{text[:4000]}..."  # Truncate to 4000 characters to fit within token limit
    
    response = bedrock_client.invoke_model(
        modelId="meta.llama3-1-70b-instruct-v1:0",
        body=json.dumps({
            "prompt": prompt
        })
    )
    
    response_body = json.loads(response['body'].read())
    return response_body['generation']



if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info", log_config=None)