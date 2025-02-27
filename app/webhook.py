from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import PlainTextResponse
from pywa import WhatsApp
from pywa.types import Message
import json
import logging
from manish import MaNish
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from typing import Optional
from fastapi.responses import JSONResponse, PlainTextResponse
import uvicorn
import PyPDF2
from pptx import Presentation
import io
import boto3
from botocore.exceptions import ClientError
import json
from pydantic import BaseModel
import httpx
from qdrant_client import QdrantClient
from qdrant_client.http import models
import uuid
import logging
import time
import requests
import base64
from dotenv import load_dotenv
import os
import aiohttp
import tempfile
from pydub import AudioSegment


load_dotenv()

def encode_image(image_file):
    return base64.b64encode(image_file.read()).decode('utf-8')


logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                    handlers=[logging.StreamHandler()])

logger = logging.getLogger("app")

app = FastAPI()
manish = MaNish(os.getenv('FACEBOOK_ACCESS_TOKEN'), phone_number_id=os.getenv('PHONE_NUMBER_ID'))

# Initialize Qdrant client
qdrant_client = QdrantClient(
    url=os.getenv('QDRANT_URL'),
    api_key=os.getenv('QDRANT_API_KEY'),
)

# Initialize AWS Bedrock client
bedrock_client = boto3.client(
    service_name="bedrock-runtime",
    region_name=os.getenv('AWS_REGION'),
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

# Initialize S3 client
s3_client = boto3.client(
    's3',
    region_name=os.getenv('AWS_REGION'),
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

def process_audio_message(data, mobile, manish):
    try:
        audio = manish.get_audio(data)
        audio_id, mime_type = audio["id"], audio["mime_type"]
        audio_url = manish.query_media_url(audio_id)
        audio_filename = manish.download_media(audio_url, mime_type)
        
        # Upload the audio file to S3
        s3_audio_url = upload_file_to_s3(audio_filename, 'meta-llama-bucket', object_name=f"audio_{audio_id}")
        
        if s3_audio_url:
            logger.info(f"{mobile} sent audio {audio_filename}, uploaded to S3: {s3_audio_url}")
        else:
            logger.error(f"Failed to upload audio {audio_filename} to S3")
        
        # Determine the sample rate (you might need to extract this from the audio file)
        sample_rate = 16000  # Example: 16kHz
        
        # Stream the audio file to Amazon Transcribe
        transcription = ""  # You would implement the transcription logic here
        
        return audio_filename, s3_audio_url, transcription
    except Exception as e:
        logger.error(f"Error processing audio message: {e}")
        return None, None, None

def upload_file_to_s3(file_path, bucket_name, object_name=None):
    """
    Upload a file to an S3 bucket and return its public URL.

    :param file_path: File to upload
    :param bucket_name: Bucket to upload to
    :param object_name: S3 object name. If not specified, file_name is used
    :return: Public URL of the uploaded file if successful, else None
    """
    # If S3 object_name was not specified, use file_name
    if object_name is None:
        object_name = os.path.basename(file_path)

    try:
        # Upload the file
        s3_client.upload_file(file_path, bucket_name, object_name)
        
        # Make the object public
        s3_client.put_object_acl(Bucket=bucket_name, Key=object_name, ACL='public-read')
        
        # Generate the public URL
        public_url = f"https://{bucket_name}.s3.amazonaws.com/{object_name}"
        return public_url
    except ClientError as e:
        print(f"Error uploading file to S3: {e}")
        return None

transcribe_client = boto3.client('transcribe', 
                                 region_name=os.getenv('AWS_REGION'), # AWS_REGION,
                                 aws_access_key_id=os.getenv('AWS_ACCESS_KEY'),
                                 aws_secret_access_key=os.getenv('AWS_SECRET_KEY'))

# Other variables
TUNE_AI_API_URL = os.getenv('TUNE_AI_API_URL')
TUNE_AI_API_KEY = os.getenv('TUNE_AI_API_KEY')

def extract_text_from_pdf(content):
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

async def convert_ogg_to_mp3(ogg_file_path):
    # Create a temporary file for the MP3
    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_mp3:
        mp3_file_path = temp_mp3.name

    # Convert OGG to MP3
    audio = AudioSegment.from_ogg(ogg_file_path)
    audio.export(mp3_file_path, format="mp3")

    return mp3_file_path

def generate_embedding(text):
    try:
        logger.info("Generating embedding")
        embedding_response = bedrock_client.invoke_model(
            modelId="amazon.titan-embed-text-v1",
            body=json.dumps({
                "inputText": text
            })
        )
        
        logger.debug(f"Raw Bedrock response: {embedding_response}")
        
        embedding_body = json.loads(embedding_response['body'].read())
        logger.debug(f"Parsed embedding body: {embedding_body}")
        
        if 'embedding' not in embedding_body:
            logger.error(f"'embedding' key not found in response. Full response: {embedding_body}")
            raise ValueError("Unexpected response format from Bedrock API")
        
        embedding = embedding_body['embedding']
        logger.info(f"Embedding generated successfully. Length: {len(embedding)}")
        return embedding
    except Exception as e:
        logger.exception(f"Error generating embedding: {str(e)}")
        raise

async def stream_audio_to_transcribe(audio_stream, sample_rate):
    # Generator function to yield audio chunks
    async def audio_stream_generator():
        async for chunk in audio_stream.iter_any():
            yield {'AudioEvent': {'AudioChunk': chunk}}

    try:
        # Start streaming transcription
        response = transcribe_client.start_stream_transcription(
            LanguageCode='en-US',
            MediaSampleRateHertz=sample_rate,
            MediaEncoding='pcm',
            AudioStream=audio_stream_generator()
        )

        # Process the streaming response
        async for event in response['TranscriptResultStream']:
            yield event

    except Exception as e:
        print(f"Error in transcription: {str(e)}")
        yield None



@app.post("/webhook", include_in_schema=False)
async def webhook(request: Request):
    data = await request.json()
    changed_field = manish.changed_field(data)
    if changed_field == "messages":
        new_message = manish.get_mobile(data)
        if new_message:
            mobile = manish.get_mobile(data)
            name = manish.get_name(data)
            message_type = manish.get_message_type(data)
            logger.info(
                f"New Message; sender:{mobile} name:{name} type:{message_type}"
            )
            if message_type == "text":
                message = manish.get_message(data)
                logger.info("Message: %s", message)
                
                try:
                    # Generate embedding for the question
                    try:
                        question_embedding = generate_embedding(message)
                    except Exception as e:
                        logger.error(f"Error generating question embedding: {str(e)}")
                        raise HTTPException(status_code=500, detail=f"Error generating question embedding: {str(e)}")

                    # Search for similar vectors in Qdrant
                    try:
                        search_results = qdrant_client.search(
                            collection_name="memories",
                            query_vector=question_embedding,
                            limit=3  # Adjust as needed
                        )
                    except Exception as e:
                        logger.exception(f"Error searching Qdrant: {str(e)}")
                        raise HTTPException(status_code=500, detail=f"Error searching vector database: {str(e)}")

                    # Retrieve and compile related content
                    context = ""
                    for result in search_results:
                        context += result.payload.get('content', '') + "\n\n"

                    # Prepare the message for Tune AI
                    messages = [
                        {
                            "role": "system",
                            "content": f"Context: {context}\nPlease provide a concise and accurate answer to the question based on the given context. If the context doesn't contain relevant information to answer the question, please state that you don't have enough information to provide an accurate answer. If the user greets , greet them back with a greeting and a pun about LaaMA Scholar whatsapp bot, also during greeting do not use any other context.."
                        },
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": message
                                }
                            ]
                        }
                    ]

                    # Prepare the request payload for Tune AI
                    payload = {
                        "temperature": 0.7,
                        "messages": messages,
                        "model": "meta/llama-3.2-90b-vision",  # Adjust if needed for text-only tasks
                        "stream": False,
                        "frequency_penalty": 0,
                        "max_tokens": 900
                    }

                    # Make the API call to Tune AI
                    try:
                        headers = {
                            "Authorization": f"Bearer {TUNE_AI_API_KEY}",
                            "Content-Type": "application/json"
                        }
                        response = requests.post(TUNE_AI_API_URL, headers=headers, json=payload)
                        response.raise_for_status()
                        result = response.json()

                        # Extract the answer from the response
                        answer = result['choices'][0]['message']['content'].strip()

                        if not answer:
                            logger.error(f"Unexpected response from Tune AI: {result}")
                            raise ValueError("Unexpected response format from Tune AI API")

                        # Send the generated answer
                        manish.send_message(f"Hi {name}, here's what I found:\n\n{answer}", mobile)

                    except requests.exceptions.RequestException as e:
                        logger.exception(f"Error calling Tune AI API: {str(e)}")
                        raise HTTPException(status_code=500, detail=f"Error generating answer: {str(e)}")

                except HTTPException as he:
                    logger.error(f"HTTP error occurred: {str(he)}")
                    manish.send_message("I'm sorry, but I encountered an error while processing your message. Please try again later.", mobile)
                except Exception as e:
                    logger.exception(f"Unexpected error in text message processing: {str(e)}")
                    manish.send_message("I apologize, but an unexpected error occurred. Please try again later.", mobile)
            
    
            elif message_type == "interactive":
                message_response = manish.get_interactive_response(data)
                intractive_type = message_response.get("type")
                message_id = message_response[intractive_type]["id"]
                message_text = message_response[intractive_type]["title"]
                logger.info(f"Interactive Message; {message_id}: {message_text}")
                manish.send_message(message_text, mobile)
            elif message_type == "location":
                message_location = manish.get_location(data)
                message_latitude = message_location["latitude"]
                message_longitude = message_location["longitude"]
                logger.info("Location: %s, %s", message_latitude, message_longitude)
                manish.send_message(f"Latitude: {message_latitude}, Longitude: {message_longitude}", mobile)
            elif message_type == "image":
                try:
                    image = manish.get_image(data)
                    image_id, mime_type = image["id"], image["mime_type"]
                    image_url = manish.query_media_url(image_id)
                    image_filename = manish.download_media(image_url, mime_type)
                    logger.info(f"{mobile} sent image {image_filename}")

                    # Read and encode the image
                    with open(image_filename, "rb") as image_file:
                        base64_image = encode_image(image_file)

                    # Prepare the message for Tune AI
                    messages = [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/jpeg;base64,{base64_image}"
                                    }
                                },
                                {
                                    "type": "text",
                                    "text": "Explain what is shown in this image from an educational point of view. Try to explain in simplest terms and give tips for exams."
                                }
                            ]
                        }
                    ]

                    # Prepare the request payload for Tune AI
                    payload = {
                        "temperature": 0.8,
                        "messages": messages,
                        "model": "meta/llama-3.2-90b-vision",
                        "stream": False,
                        "frequency_penalty": 0,
                        "max_tokens": 900
                    }

                    # Make the API call to Tune AI
                    headers = {
                        "Authorization": f"Bearer {TUNE_AI_API_KEY}",
                        "Content-Type": "application/json"
                    }
                    response = requests.post(TUNE_AI_API_URL, headers=headers, json=payload)
                    response.raise_for_status()
                    result = response.json()

                    # Extract the analysis result from the response
                    analysis_result = result['choices'][0]['message']['content']

                    # Send the analysis result back to the user
                    manish.send_message(f"Here's what I see in the image:\n\n{analysis_result}", mobile)

                except Exception as e:
                    logger.error(f"Error processing image: {str(e)}")
                    manish.send_message("I'm sorry, but I encountered an error while processing your image. Please try again later.", mobile)
            elif message_type == "video":
                video = manish.get_video(data)
                video_id, mime_type = video["id"], video["mime_type"]
                video_url = manish.query_media_url(video_id)
                video_filename = manish.download_media(video_url, mime_type)
                logger.info(f"{mobile} sent video {video_filename}")
                manish.send_message(f"Video: {video_filename}", mobile)
            elif message_type == "audio":
                audio_filename, s3_audio_url, transcription = process_audio_message(data, mobile, manish)
                if audio_filename and s3_audio_url:
                    # Process the transcription or perform other tasks
                    if transcription:
                        manish.send_message(f"Here's the transcription of your audio message:\n\n{transcription.strip()}", mobile)
                    else:
                        manish.send_message("I'm sorry, but I couldn't transcribe your audio message. It might be too short or unclear.", mobile)
                else:
                    logger.error("Failed to process audio message")
            elif message_type == "file":
                file = manish.get_file(data)
                logger.info(file)
                file_id, mime_type = file["id"], file["mime_type"]
                file_url = manish.query_media_url(file_id)
                file_filename = manish.download_media(file_url, mime_type)
                logger.info(f"{mobile} sent file {file_filename}")
                manish.send_message(f"File: {file_filename}", mobile)
            elif message_type == "document":
                document = manish.get_document(data)
                document_id, mime_type = document["id"], document["mime_type"]
                document_url = manish.query_media_url(document_id)
                document_filename = manish.download_media(document_url, mime_type)
                logger.info(f"{mobile} sent document {document_filename}")

                # Process the document similar to the /remember endpoint
                try:
                    # Read the document content
                    with open(document_filename, 'rb') as file:
                        content = file.read()

                    # Extract text based on file type
                    if document_filename.endswith('.pdf'):
                        text = extract_text_from_pdf(content)
                    elif document_filename.endswith('.pptx'):
                        text = extract_text_from_pptx(content)
                    else:
                        raise ValueError("Unsupported file format. Please upload a PDF or PPTX file.")

                    # Generate embedding
                    embedding_response = bedrock_client.invoke_model(
                        modelId="amazon.titan-embed-text-v1",
                        body=json.dumps({"inputText": text})
                    )
                    embedding_body = json.loads(embedding_response['body'].read())
                    embedding = embedding_body['embedding']

                    # Generate a unique ID for this memory
                    memory_id = str(uuid.uuid4())

                    # Save to Qdrant
                    qdrant_client.upsert(
                        collection_name="memories",
                        points=[
                            models.PointStruct(
                                id=memory_id,
                                vector=embedding,
                                payload={
                                    "content": text,
                                    "filename": document_filename,
                                    "sender": mobile,
                                    "name": name
                                }
                            )
                        ]
                    )

                    manish.send_message(f"Document processed and stored successfully. Memory ID: {memory_id}", mobile)
                except Exception as e:
                    logger.error(f"Error processing document: {str(e)}")
                    manish.send_message(f"Error processing document: {str(e)}", mobile)

            else:
                logger.info(f"{mobile} sent {message_type} ")
                manish.send_message(f"Sorry , I don't know how to handle {message_type}", mobile)
                logger.info(data)
        else:
            delivery = manish.get_delivery(data)
            if delivery:
                logger.info(f"Message : {delivery}")
            else:
                logger.info("No new message")
    return "ok"


VERIFY_TOKEN = "meatyhamhock"

@app.get("/webhook", include_in_schema=False)
async def verify(request: Request):
    if request.query_params.get('hub.mode') == "subscribe" and request.query_params.get("hub.challenge"):
        if not request.query_params.get('hub.verify_token') == VERIFY_TOKEN: #os.environ["VERIFY_TOKEN"]:
           return "Verification token mismatch", 403
        return int(request.query_params.get('hub.challenge'))
    return "Hello world", 200

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
